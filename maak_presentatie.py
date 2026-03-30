from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

SS = '/home/user/demo/screenshots/'
OUT = '/home/user/demo/vth-export/OmgevingsCheck_Presentatie.pptx'

# Kleuren
ORANJE   = RGBColor(0xC4, 0x62, 0x2D)
GROEN    = RGBColor(0x1B, 0x6B, 0x6B)
WIT      = RGBColor(0xFF, 0xFF, 0xFF)
LICHTGRIJS = RGBColor(0xF5, 0xF5, 0xF5)
DONKERBLAUW = RGBColor(0x1A, 0x23, 0x3A)
TEKSTGRIJS = RGBColor(0x55, 0x65, 0x6E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # volledig leeg

def add_rect(slide, l, t, w, h, fill=None, line=None, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape

def add_text(slide, txt, l, t, w, h, size=18, bold=False, color=WIT, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box

def add_img(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def header_band(slide, titel, subtitel=None):
    """Groene balk bovenaan"""
    add_rect(slide, 0, 0, 13.33, 1.4, fill=GROEN)
    add_text(slide, titel, 0.4, 0.15, 10, 0.75, size=30, bold=True, color=WIT)
    if subtitel:
        add_text(slide, subtitel, 0.4, 0.85, 10, 0.45, size=14, color=RGBColor(0xCC, 0xE8, 0xE8))
    # Oranje accent balk
    add_rect(slide, 0, 1.4, 13.33, 0.06, fill=ORANJE)

def footer(slide):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=DONKERBLAUW)
    add_text(slide, 'OmgevingsCheck App  ·  Gemeente Waalre  ·  © 2026 Christian Krielen', 0.4, 7.2, 12, 0.3,
             size=9, color=RGBColor(0xAA, 0xBB, 0xCC))

def bullet_slide(slide, items, l=0.5, t=1.7, size=15, color=DONKERBLAUW, symbol='●'):
    y = t
    for item in items:
        if isinstance(item, tuple):
            sym, tekst = item
        else:
            sym, tekst = symbol, item
        add_text(slide, sym, l, y, 0.35, 0.38, size=size, color=ORANJE, bold=True)
        add_text(slide, tekst, l+0.35, y, 12.0, 0.38, size=size, color=color)
        y += 0.40
    return y

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITELPAGINA
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DONKERBLAUW)
add_rect(sl, 0, 0, 13.33, 4.2, fill=GROEN)
# Decoratief oranje vlak
add_rect(sl, 0, 4.2, 13.33, 0.12, fill=ORANJE)
# App naam
add_text(sl, 'OmgevingsCheck', 0.6, 0.6, 12, 1.2, size=56, bold=True, color=WIT, align=PP_ALIGN.LEFT)
add_text(sl, 'Digitaal bouwtoezicht voor gemeente Waalre', 0.6, 1.8, 11, 0.7, size=22, color=RGBColor(0xCC,0xE8,0xE8))
add_text(sl, 'Ontwikkeld door: Christian Krielen  |  Afdeling VTH', 0.6, 2.55, 10, 0.5, size=16, color=RGBColor(0xAA,0xCC,0xCC))
add_text(sl, 'Presentatie t.b.v. management  —  2026', 0.6, 3.1, 10, 0.5, size=14, color=RGBColor(0x88,0xAA,0xAA))
# Mockup screenshot rechts
add_img(sl, SS+'01_projectlijst.png', 8.8, 0.3, 4.1, 3.7)
# Ondertekst
add_text(sl, '🔒  Volledig offline  ·  📱  Tablet & smartphone  ·  🏛️  AVG-compliant', 0.6, 4.6, 12, 0.6,
         size=15, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.LEFT)
add_text(sl, '© 2026 Christian Krielen  —  Gemeente Waalre', 0.6, 7.1, 12, 0.3, size=10, color=RGBColor(0x77,0x88,0x99))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — AANLEIDING & CONTEXT
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'Aanleiding', 'Waarom is OmgevingsCheck ontwikkeld?')
footer(sl)

add_text(sl, 'Het probleem', 0.5, 1.65, 6, 0.4, size=16, bold=True, color=GROEN)
bullet_slide(sl, [
    'Toezichthouders werkten met papieren formulieren op locatie',
    'Controleresultaten werden achteraf handmatig overgetypt',
    'Foto\'s werden los opgeslagen, moeilijk te koppelen aan dossiers',
    'Geen eenvoudig overzicht van de status per bouwproject',
    'Rapportages kostten veel tijd en waren foutgevoelig',
], t=2.05, size=14, color=DONKERBLAUW)

add_text(sl, 'De wens', 0.5, 4.35, 6, 0.4, size=16, bold=True, color=GROEN)
bullet_slide(sl, [
    'Eén digitaal hulpmiddel, altijd bij de hand op tablet',
    'Directe registratie op locatie, ook zonder internet',
    'Gestructureerde controles gekoppeld aan BBL-artikelen',
    'Snel en professioneel rapporteren',
], t=4.75, size=14, color=DONKERBLAUW)

add_img(sl, SS+'01_projectlijst.png', 7.4, 1.55, 5.5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — WAT IS OMGEVINGSCHECK?
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'Wat is OmgevingsCheck?', 'Een offline webapp voor gemeentelijk bouwtoezicht')
footer(sl)

add_img(sl, SS+'desktop_01_projectlijst.png', 0.3, 1.55, 8.0)
add_text(sl, 'Kernkenmerken', 8.5, 1.55, 4.5, 0.4, size=16, bold=True, color=GROEN)
bullet_slide(sl, [
    'Werkt 100% offline',
    'Draait in elke browser',
    'Tablet & smartphone',
    'Geen installatie nodig',
    'Data blijft lokaal',
    'Koppeling PDOK / kadaster',
    'Import uit zakensysteem',
    'PDF/e-mail rapportage',
], l=8.5, t=2.0, size=13, color=DONKERBLAUW)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — FUNCTIONALITEITEN OVERZICHT
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'Functionaliteiten', 'Alles wat een toezichthouder nodig heeft')
footer(sl)

# 6 tegels
tiles = [
    ('📋', 'Projectbeheer', 'Zaken aanmaken, bewerken\nen filteren op status'),
    ('🗺️', 'Kaartweergave', 'Alle bouwprojecten op\ninteractieve kaart (PDOK)'),
    ('🔍', 'Controles', 'Gestructureerde controles\nmet BBL-artikelkoppeling'),
    ('📊', 'Voortgang', 'Overzicht naleving van\nvergunningvoorschriften'),
    ('📥', 'XLS-import', 'Zaken direct importeren\nuit het zakensysteem'),
    ('📧', 'Rapportage', 'Rapport per e-mail sturen\nnaar Outlook'),
]
cols = [(0.25, 4.3), (4.55, 4.3), (8.85, 4.3)]
rows = [1.6, 4.2]
ti = 0
for row in rows:
    for col, cw in cols:
        if ti >= len(tiles): break
        icon, titel, txt = tiles[ti]
        add_rect(sl, col, row, cw, 2.35, fill=WIT)
        add_text(sl, icon, col+0.15, row+0.15, 0.7, 0.6, size=26, color=GROEN)
        add_text(sl, titel, col+0.85, row+0.2, 3.2, 0.5, size=15, bold=True, color=DONKERBLAUW)
        add_text(sl, txt, col+0.15, row+0.75, 3.9, 1.4, size=12, color=TEKSTGRIJS, wrap=True)
        ti += 1

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — PROJECTLIJST & KAART
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'Projectenoverzicht & Kaart', 'Filteren, zoeken en locaties visualiseren')
footer(sl)

add_img(sl, SS+'desktop_01_projectlijst.png', 0.2, 1.55, 6.4)
add_img(sl, SS+'desktop_02_kaart.png', 6.75, 1.55, 6.35)

add_text(sl, '← Lijst met filters op status / zaaktype', 0.2, 6.65, 6, 0.4, size=11, color=TEKSTGRIJS)
add_text(sl, '← Kaartweergave met gekleurde markers per status', 6.75, 6.65, 6, 0.4, size=11, color=TEKSTGRIJS)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — PROJECTDETAIL & CONTROLES
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'Projectdetail & Controles', 'Van basisgegevens tot controlerapport')
footer(sl)

add_img(sl, SS+'desktop_03_detail.png', 0.2, 1.55, 4.2)
add_img(sl, SS+'desktop_04_controles.png', 4.55, 1.55, 4.2)
add_img(sl, SS+'desktop_05_voortgang.png', 8.9, 1.55, 4.2)

add_text(sl, 'Basisgegevens', 0.2, 6.68, 4, 0.35, size=11, color=TEKSTGRIJS, align=PP_ALIGN.CENTER)
add_text(sl, 'Controles per datum', 4.55, 6.68, 4, 0.35, size=11, color=TEKSTGRIJS, align=PP_ALIGN.CENTER)
add_text(sl, 'Voortgang voorschriften', 8.9, 6.68, 4, 0.35, size=11, color=TEKSTGRIJS, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — BBL ARTIKELKOPPELING
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'BBL-artikelkoppeling', 'Controles gebaseerd op wettelijke grondslag')
footer(sl)

add_img(sl, SS+'09_controle_detail.png', 0.3, 1.55, 3.5)
add_text(sl, 'Ingebouwde juridische bibliotheek', 4.2, 1.65, 8.8, 0.5, size=18, bold=True, color=GROEN)
bullet_slide(sl, [
    '352 artikelen uit het Besluit Bouwwerken Leefomgeving (BBL)',
    'Hoofdstuk 3: Veiligheid bij brand + Gezondheid',
    'Hoofdstuk 4: Veiligheid, Gezondheid, Bouwwerkinstallaties',
    'Elk controleonderwerp koppelt aan een specifiek BBL-artikel',
    'Zoeken op artikelnummer, trefwoord of onderwerp',
    'Directe link naar officiële wetgeving (wetten.overheid.nl)',
    'Juridisch onderbouwde rapportages, klaar voor handhaving',
], l=4.2, t=2.25, size=13, color=DONKERBLAUW)

add_text(sl, '← Controledetail op tablet', 0.3, 6.7, 3.5, 0.35, size=11, color=TEKSTGRIJS)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — XLS IMPORT
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'Koppeling met het zakensysteem', 'Zaken importeren via XLS-export')
footer(sl)

add_img(sl, SS+'desktop_06_xls_import.png', 0.2, 1.55, 7.8)
add_text(sl, 'Hoe werkt het?', 8.3, 1.65, 4.8, 0.45, size=16, bold=True, color=GROEN)
bullet_slide(sl, [
    ('1.', 'XLS-export uit zakensysteem'),
    ('2.', 'Importeer in OmgevingsCheck'),
    ('3.', 'Preview met checkboxes'),
    ('4.', 'Duplicaatdetectie automatisch'),
    ('5.', 'Adressen geocodeerd via PDOK'),
    ('6.', 'Direct zichtbaar op kaart'),
], l=8.3, t=2.2, size=13, color=DONKERBLAUW, symbol='')

add_text(sl, '← Importscherm met preview en selectie', 0.2, 6.7, 7.8, 0.35, size=11, color=TEKSTGRIJS)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — AVG & PRIVACY
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'AVG & Privacy', 'Offline-first architectuur beschermt persoonsgegevens')
footer(sl)

# Groen kader: wat NIET het internet op gaat
add_rect(sl, 0.3, 1.6, 6.0, 4.9, fill=RGBColor(0xE8,0xF5,0xE9))
add_text(sl, '🔒  Blijft lokaal op het apparaat', 0.5, 1.7, 5.5, 0.45, size=15, bold=True, color=GROEN)
bullet_slide(sl, [
    'Projectgegevens (naam, adres, status)',
    'Controleverslagen en bevindingen',
    'Foto\'s van de bouwplaats',
    'Vergunningsdocumenten (PDF)',
    'Alle gebruikersinstellingen',
    'Backup-bestanden',
], l=0.5, t=2.25, size=13, color=DONKERBLAUW)

# Oranje kader: wat WEL het internet op gaat
add_rect(sl, 6.7, 1.6, 6.3, 4.9, fill=RGBColor(0xFFF3,0xE0,0x00)[:3] if False else RGBColor(0xFF,0xF3,0xE0))
add_text(sl, '🌐  Externe diensten (functioneel)', 6.9, 1.7, 5.8, 0.45, size=15, bold=True, color=ORANJE)
items2 = [
    ('Kaartafbeeldingen', 'OpenStreetMap tiles\n(geen persoonsdata)'),
    ('Adresgeocodering', 'PDOK Locatieserver\n(alleen adres, geen naam)'),
    ('Kaartweergave', 'Leaflet.js bibliotheek\n(software, geen data)'),
    ('E-mail', 'Via Outlook op tablet\n(gebruiker verstuurt zelf)'),
]
y = 2.25
for dienst, uitleg in items2:
    add_text(sl, '▸  ' + dienst, 6.9, y, 2.8, 0.32, size=13, bold=True, color=DONKERBLAUW)
    add_text(sl, uitleg, 9.75, y, 3.0, 0.35, size=11, color=TEKSTGRIJS)
    y += 0.62

add_text(sl, '✅  Alle persoonsgebonden projectdata verlaat nooit het apparaat — volledig AVG-compliant',
         0.3, 6.6, 12.7, 0.45, size=12, bold=True, color=GROEN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — EFFICIENCY
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'Efficiency & Tijdwinst', 'Minder administratie, meer toezicht')
footer(sl)

# 3 kolommen
for i, (icon, titel, voor, na) in enumerate([
    ('⏱️', 'Rapportage opstellen',
     'Handmatig typen:\n30–60 minuten per controle',
     'OmgevingsCheck:\n5 minuten, direct verzenden'),
    ('📁', 'Dossierbeheer',
     'Papier + losse bestanden:\nmoeilijk doorzoekbaar',
     'Alles digitaal gekoppeld:\nin 2 tikken gevonden'),
    ('📍', 'Locatie opzoeken',
     'Kaart raadplegen of\nbellen voor adres',
     'PDOK geocodering:\nautomatisch op de kaart'),
]):
    x = 0.3 + i * 4.3
    add_rect(sl, x, 1.6, 4.1, 5.4, fill=WIT)
    add_text(sl, icon, x+0.2, 1.75, 0.7, 0.6, size=28)
    add_text(sl, titel, x+0.9, 1.8, 3.0, 0.55, size=15, bold=True, color=DONKERBLAUW)
    add_rect(sl, x+0.2, 2.55, 3.7, 1.6, fill=RGBColor(0xFF,0xEE,0xE8))
    add_text(sl, '❌  Vroeger', x+0.3, 2.6, 3.5, 0.35, size=11, bold=True, color=ORANJE)
    add_text(sl, voor, x+0.3, 2.95, 3.5, 1.1, size=11, color=TEKSTGRIJS, wrap=True)
    add_rect(sl, x+0.2, 4.3, 3.7, 1.8, fill=RGBColor(0xE8,0xF5,0xE9))
    add_text(sl, '✅  Nu', x+0.3, 4.35, 3.5, 0.35, size=11, bold=True, color=GROEN)
    add_text(sl, na, x+0.3, 4.7, 3.5, 1.1, size=11, color=DONKERBLAUW, wrap=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — TOEKOMST & KANSEN
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LICHTGRIJS)
header_band(sl, 'Kansen & Toekomst', 'Doorontwikkeling en bredere inzet')
footer(sl)

kansen = [
    ('🏛️', 'Andere gemeenten',     'OmgevingsCheck is generiek inzetbaar voor elke VTH-organisatie in Nederland'),
    ('👥', 'Meerdere inspecteurs', 'Uitbreiding naar teamgebruik met gedeelde projectdossiers via een centrale server'),
    ('🔗', 'Systeemkoppeling',     'Directe API-koppeling met zakensysteem (bijv. Rx.Mission) voor automatische sync'),
    ('📱', 'Native app',           'Verpakken als Progressive Web App (PWA) of native iOS/Android-app'),
    ('🤖', 'Slimme controles',     'AI-ondersteuning bij het bepalen relevante BBL-artikelen o.b.v. vergunningstype'),
    ('📊', 'Managementrapportage', 'Dashboards met statistieken: doorlooptijden, hercontroles, naleving per wijk'),
]
y = 1.65
for i, (icon, titel, txt) in enumerate(kansen):
    col = 0 if i % 2 == 0 else 6.7
    row = y + (i // 2) * 1.7
    add_rect(sl, col+0.2, row, 6.2, 1.5, fill=WIT)
    add_text(sl, icon, col+0.35, row+0.15, 0.7, 0.7, size=26)
    add_text(sl, titel, col+1.1, row+0.18, 4.8, 0.45, size=14, bold=True, color=GROEN)
    add_text(sl, txt, col+1.1, row+0.62, 5.0, 0.75, size=11, color=TEKSTGRIJS, wrap=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — SAMENVATTING & VOLGENDE STAP
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DONKERBLAUW)
add_rect(sl, 0, 0, 13.33, 2.0, fill=GROEN)
add_rect(sl, 0, 2.0, 13.33, 0.1, fill=ORANJE)

add_text(sl, 'Samenvatting', 0.7, 0.2, 12, 0.75, size=34, bold=True, color=WIT)
add_text(sl, 'OmgevingsCheck is klaar voor dagelijks gebruik', 0.7, 0.9, 12, 0.6, size=17, color=RGBColor(0xCC,0xE8,0xE8))

samenvatting = [
    '✅  Volledig werkende offline webapp voor VTH bouwtoezicht',
    '✅  352 BBL-artikelen (H3 + H4) ingebouwd als juridische bibliotheek',
    '✅  Kaartweergave met PDOK-geocodering voor alle bouwprojecten',
    '✅  XLS-import direct vanuit het zakensysteem',
    '✅  Rapportage via e-mail (Outlook) met één druk op de knop',
    '✅  Volledig AVG-compliant — persoonsdata verlaat het apparaat nooit',
    '✅  Backup & herstel inclusief foto\'s en documenten',
]
y = 2.25
for item in samenvatting:
    add_text(sl, item, 0.7, y, 12, 0.38, size=14, color=WIT)
    y += 0.42

add_rect(sl, 0.7, 5.8, 5.5, 1.35, fill=RGBColor(0x22,0x3A,0x5C))
add_text(sl, '📬  Volgende stap', 0.9, 5.85, 5.0, 0.4, size=14, bold=True, color=ORANJE)
add_text(sl, 'Formele ingebruikname + evaluatie na 3 maanden\nEventuele aansluiting andere inspecteurs bespreken', 0.9, 6.25, 5.0, 0.8, size=12, color=WIT, wrap=True)

add_rect(sl, 6.9, 5.8, 6.1, 1.35, fill=RGBColor(0x22,0x3A,0x5C))
add_text(sl, '💬  Vragen?', 7.1, 5.85, 5.6, 0.4, size=14, bold=True, color=ORANJE)
add_text(sl, 'Christian Krielen\nckrielen@waalre.nl', 7.1, 6.25, 5.6, 0.8, size=12, color=WIT)

add_text(sl, '© 2026 Christian Krielen  —  Gemeente Waalre', 0.7, 7.25, 12, 0.2, size=9, color=RGBColor(0x77,0x88,0x99))

# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'✅ Presentatie opgeslagen: {OUT}')
print(f'   Slides: {len(prs.slides)}')
